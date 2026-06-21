"""
Agent Tools - Built-in capabilities for the AI Agent.
Includes web_fetch_batch, image_download, create_ppt, etc.
web_search and image_search are now provided via MCP protocol.
"""

import os
import json
import re
import requests
from typing import List, Optional, Dict, Any
from urllib.parse import urlparse
from datetime import datetime

from modules.config.config_manager import ConfigManager
from modules.mcp.mcp_registry import registry as mcp_registry
from modules.utils.logger import log_agent_action, log_error, log_info, log_tool_call


class AgentTools:
    """Collection of all agent capabilities."""

    def __init__(self):
        self.config = ConfigManager()
        self._download_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "download")
        self._output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "output")
        os.makedirs(self._download_dir, exist_ok=True)
        os.makedirs(self._output_dir, exist_ok=True)

    def web_search_func(self, query: str) -> List[Dict[str, str]]:
        """Execute web search via MCP (fallback to direct)."""
        log_tool_call(f"web_search(query='{query}') via MCP")
        try:
            result_text = mcp_registry.call_tool("web_search", {"query": query})
            if result_text:
                results = json.loads(result_text)
                if isinstance(results, list):
                    return results
            return []
        except Exception:
            # Fallback: direct SearXNG call if MCP unavailable
            from modules.utils.logger import log_warning
            log_warning(f"MCP web_search unavailable, using direct fallback")
            return self._fallback_web_search(query)

    def _fallback_web_search(self, query: str, max_results: int = 10) -> List[Dict[str, str]]:
        """Direct SearXNG call as fallback when MCP is not available."""
        base_url = self.config.get("tools.searxng.base_url", "http://localhost:8888")
        try:
            resp = requests.post(
                f"{base_url}/search",
                data={
                    "q": query, "format": "json", "language": "zh-CN",
                    "categories": "general", "pageno": 1,
                },
                timeout=15,
                headers={"Accept": "application/json"}
            )
            if resp.status_code == 200:
                data = resp.json()
                results = []
                for item in data.get("results", [])[:max_results]:
                    results.append({
                        "title": item.get("title", ""),
                        "url": item.get("url", ""),
                        "content": item.get("content", ""),
                    })
                return results
        except Exception:
            pass
        return [{
            "title": f"Mock Result: {query} - Overview",
            "url": f"https://example.com/result?q={query.replace(' ', '+')}",
            "content": f"Mock search result for '{query}'. Configure MCP SearXNG server for real results."
        }]

    def image_search_func(self, query: str, max_results: int = 5) -> List[Dict[str, str]]:
        """Search for images via MCP (fallback to direct)."""
        log_tool_call(f"image_search(query='{query}', max_results={max_results}) via MCP")
        try:
            result_text = mcp_registry.call_tool("image_search", {"query": query, "max_results": max_results})
            if result_text:
                results = json.loads(result_text)
                if isinstance(results, list):
                    return results
            return []
        except Exception:
            from modules.utils.logger import log_warning
            log_warning(f"MCP image_search unavailable, using direct fallback")
            return self._fallback_image_search(query, max_results)

    def _fallback_image_search(self, query: str, max_results: int = 5) -> List[Dict[str, str]]:
        """Direct image API call as fallback when MCP is not available."""
        config = self.config.get("tools.image_search", {})
        provider = config.get("provider", "pexels")
        if provider == "pexels":
            api_key = config.get("pexels.api_key", "") or config.get("pexels", {}).get("api_key", "")
            if api_key:
                try:
                    resp = requests.get(
                        "https://api.pexels.com/v1/search",
                        params={"query": query, "per_page": max_results},
                        headers={"Authorization": api_key}, timeout=10
                    )
                    if resp.status_code == 200:
                        return [{"url": p.get("src", {}).get("original", ""),
                                 "small_url": p.get("src", {}).get("medium", ""),
                                 "photographer": p.get("photographer", ""),
                                 "source": "pexels", "alt": p.get("alt", "")}
                                for p in resp.json().get("photos", [])]
                except Exception:
                    pass
        elif provider == "unsplash":
            api_key = config.get("unsplash.api_key", "") or config.get("unsplash", {}).get("api_key", "")
            if api_key:
                try:
                    resp = requests.get(
                        "https://api.unsplash.com/search/photos",
                        params={"query": query, "per_page": max_results},
                        headers={"Authorization": f"Client-ID {api_key}"}, timeout=10
                    )
                    if resp.status_code == 200:
                        return [{"url": p.get("urls", {}).get("regular", ""),
                                 "small_url": p.get("urls", {}).get("small", ""),
                                 "photographer": p.get("user", {}).get("name", ""),
                                 "source": "unsplash", "alt": p.get("alt_description", "")}
                                for p in resp.json().get("results", [])]
                except Exception:
                    pass
        return [{
            "url": f"https://via.placeholder.com/800x600.png?text={query.replace(' ', '+')}+1",
            "small_url": f"https://via.placeholder.com/400x300.png?text={query.replace(' ', '+')}+1",
            "photographer": "Mock", "source": "mock", "alt": query,
        }]

    def web_fetch_batch(self, urls: List[str]) -> str:
        """Fetch content from multiple URLs and combine."""
        log_tool_call(f"web_fetch_batch({len(urls)} URLs)")
        combined = []
        for i, url in enumerate(urls):
            try:
                resp = requests.get(url, timeout=15, headers={
                    "User-Agent": "Mozilla/5.0 (compatible; AI-Agent/1.0)"
                })
                if resp.status_code == 200:
                    content_type = resp.headers.get("Content-Type", "")
                    if "text/html" in content_type or "text/plain" in content_type or "application/json" in content_type:
                        text = self._extract_text(resp.text)
                        combined.append(f"=== URL [{i+1}/{len(urls)}]: {url} ===\n{text[:3000]}\n")
                    else:
                        combined.append(f"=== URL [{i+1}/{len(urls)}]: {url} ===\n[Non-text content: {content_type}, {len(resp.content)} bytes]\n")
                else:
                    combined.append(f"=== URL [{i+1}/{len(urls)}]: {url} ===\n[HTTP {resp.status_code}]\n")
                log_agent_action(f"Fetched: {url}")
            except Exception as e:
                combined.append(f"=== URL [{i+1}/{len(urls)}]: {url} ===\n[Error: {e}]\n")
                log_error(f"Failed to fetch {url}: {e}")
        return "\n".join(combined)

    def _extract_text(self, html: str) -> str:
        """Basic HTML to text extraction."""
        # Remove scripts and styles
        html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
        html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL)
        # Remove tags
        text = re.sub(r'<[^>]+>', ' ', html)
        # Decode entities
        text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
        text = text.replace('&quot;', '"').replace('&#39;', "'")
        text = text.replace('&nbsp;', ' ')
        # Clean whitespace
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def image_download(self, urls: List[str]) -> List[str]:
        """Download images from URLs to download/ directory."""
        log_tool_call(f"image_download({len(urls)} images)")
        saved = []
        for url in urls:
            try:
                resp = requests.get(url, timeout=30, stream=True)
                if resp.status_code == 200:
                    # Determine extension from URL or content-type
                    ext = ".jpg"
                    ct = resp.headers.get("Content-Type", "")
                    if "png" in ct:
                        ext = ".png"
                    elif "gif" in ct:
                        ext = ".gif"
                    elif "webp" in ct:
                        ext = ".webp"

                    # Generate unique filename
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:20]
                    safe_name = f"img_{timestamp}{ext}"
                    filepath = os.path.join(self._download_dir, safe_name)

                    with open(filepath, "wb") as f:
                        for chunk in resp.iter_content(chunk_size=8192):
                            if chunk:
                                f.write(chunk)
                    saved.append(filepath)
                    log_agent_action(f"Downloaded image: {filepath}")
                else:
                    log_error(f"Failed to download {url}: HTTP {resp.status_code}")
            except Exception as e:
                log_error(f"Failed to download {url}: {e}")
        return saved

    def create_ppt(self, config: Dict[str, Any]) -> str:
        """Generate a PowerPoint file using python-pptx."""
        log_tool_call("create_ppt()")
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            slides_data = config.get("slides", [])
            color_scheme = config.get("color_scheme", "modern_blue")

            # Predefined color schemes
            schemes = {
                "modern_blue": {
                    "primary": RGBColor(0x1A, 0x3C, 0x6E),
                    "secondary": RGBColor(0x2D, 0x7D, 0xD2),
                    "accent": RGBColor(0xFF, 0x8C, 0x00),
                    "light": RGBColor(0xF0, 0xF4, 0xFA),
                    "dark": RGBColor(0x0D, 0x1B, 0x2A),
                    "text": RGBColor(0x33, 0x33, 0x33),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "corporate_gray": {
                    "primary": RGBColor(0x2C, 0x3E, 0x50),
                    "secondary": RGBColor(0x7F, 0x8C, 0x8D),
                    "accent": RGBColor(0x34, 0x98, 0xDB),
                    "light": RGBColor(0xEC, 0xF0, 0xF1),
                    "dark": RGBColor(0x1A, 0x25, 0x2F),
                    "text": RGBColor(0x2C, 0x3E, 0x50),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "elegant_green": {
                    "primary": RGBColor(0x1E, 0x56, 0x36),
                    "secondary": RGBColor(0x4C, 0xAF, 0x50),
                    "accent": RGBColor(0xFF, 0xD7, 0x00),
                    "light": RGBColor(0xE8, 0xF5, 0xE9),
                    "dark": RGBColor(0x0D, 0x2B, 0x1A),
                    "text": RGBColor(0x33, 0x33, 0x33),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "dark_elegant": {
                    "primary": RGBColor(0x0D, 0x0D, 0x0D),
                    "secondary": RGBColor(0x1A, 0x1A, 0x2E),
                    "accent": RGBColor(0xC9, 0xA9, 0x4C),
                    "light": RGBColor(0x2A, 0x2A, 0x3E),
                    "dark": RGBColor(0x00, 0x00, 0x00),
                    "text": RGBColor(0xE0, 0xE0, 0xE0),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "tech": {
                    "primary": RGBColor(0x0A, 0x0A, 0x1A),
                    "secondary": RGBColor(0x00, 0xD4, 0xFF),
                    "accent": RGBColor(0xBB, 0x86, 0xFC),
                    "light": RGBColor(0x1A, 0x1A, 0x3A),
                    "dark": RGBColor(0x00, 0x00, 0x00),
                    "text": RGBColor(0xE0, 0xE0, 0xE0),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
            }

            scheme = schemes.get(color_scheme, schemes["modern_blue"])

            for slide_data in slides_data:
                slide_type = slide_data.get("type", "content")
                title = slide_data.get("title", "")
                content = slide_data.get("content", [])
                bg_config = slide_data.get("background", {})

                # Choose layout based on type
                if slide_type == "title_slide":
                    layout = prs.slide_layouts[6]  # Blank
                    slide = prs.slides.add_slide(layout)
                    self._apply_background(slide, bg_config, scheme)

                    # Title
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(44)
                    p.font.color.rgb = scheme["white"] if "dark" in color_scheme or "tech" in color_scheme else scheme["primary"]
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                    # Subtitle
                    subtitle = slide_data.get("subtitle", "")
                    if subtitle:
                        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.333), Inches(1))
                        tf2 = txBox2.text_frame
                        p2 = tf2.paragraphs[0]
                        p2.text = subtitle
                        p2.font.size = Pt(24)
                        p2.font.color.rgb = scheme["accent"]
                        p2.alignment = PP_ALIGN.CENTER

                elif slide_type == "section_header":
                    layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(layout)
                    self._apply_background(slide, bg_config, scheme)

                    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.5))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(40)
                    p.font.color.rgb = scheme["accent"]
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                else:
                    # Content slide
                    layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(layout)
                    self._apply_background(slide, bg_config, scheme)

                    # Title bar
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = scheme["primary"]
                    shape.line.fill.background()

                    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.9))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(32)
                    p.font.color.rgb = scheme["white"]
                    p.font.bold = True

                    # Content
                    if content:
                        left_margin = 0.8 if slide_data.get("layout") != "two_column" else 0.5
                        txBox2 = slide.shapes.add_textbox(
                            Inches(left_margin), Inches(1.5),
                            Inches(11.733 if slide_data.get("layout") != "two_column" else 5.5),
                            Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.word_wrap = True

                        for i, item in enumerate(content):
                            if i == 0:
                                p = tf2.paragraphs[0]
                            else:
                                p = tf2.add_paragraph()
                            p.text = item
                            p.font.size = Pt(18)
                            p.font.color.rgb = scheme["text"]
                            p.space_after = Pt(8)
                            p.level = 0
                            if item.startswith("  "):
                                p.level = 1
                                p.font.size = Pt(16)

                # Speaker notes
                notes = slide_data.get("notes", "")
                if notes and slide.has_notes_slide:
                    slide.notes_slide.notes_text_frame.text = notes

            # Save
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
            filepath = os.path.join(self._output_dir, filename)
            prs.save(filepath)
            log_agent_action(f"PPT saved: {filepath}")
            return filepath

        except Exception as e:
            log_error(f"PPT generation failed: {e}")
            raise

    def _apply_background(self, slide, bg_config: Dict[str, Any], scheme: Dict):
        """Apply background to a slide."""
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
        from lxml import etree

        bg_type = bg_config.get("type", "solid")
        color_1 = bg_config.get("color_1")
        color_2 = bg_config.get("color_2")

        def _to_rgb(c):
            if isinstance(c, list) and len(c) >= 3:
                return RGBColor(*c[:3])
            return scheme.get("light", RGBColor(240, 244, 250))

        c1 = _to_rgb(color_1) if color_1 else scheme.get("primary", RGBColor(26, 60, 110))
        c2 = _to_rgb(color_2) if color_2 else scheme.get("secondary", RGBColor(45, 125, 210))

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = c1

        if bg_type == "gradient":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = c2
            shape.line.fill.background()

            spPr = shape._element.find(qn('p:spPr'))
            if spPr is not None:
                solidFill_elem = spPr.find(qn('a:solidFill'))
                if solidFill_elem is not None:
                    srgbClr = solidFill_elem.find(qn('a:srgbClr'))
                    if srgbClr is not None:
                        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                        alpha.set('val', '50000')

    def save_code(self, code: str, filename: str, language: str = "py") -> str:
        """Save generated code to output directory."""
        os.makedirs(self._output_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = f"code_{timestamp}_{filename}"
        filepath = os.path.join(self._output_dir, safe_name)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(code)
        log_agent_action(f"Code saved: {filepath}")
        return filepath

    def run_code(self, filepath: str, args: str = "") -> str:
        """Run a code file and return output."""
        from modules.utils.file_ops import FileOps
        log_agent_action(f"Running code: {filepath} {args}")
        return FileOps.bash(f"python3 {filepath} {args}", timeout=30)
