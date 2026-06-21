"""
PPT Tools Plugin - PowerPoint presentation generation.
"""

import os
from typing import Any, Dict
from datetime import datetime

from modules.agents.tool_base import BaseTool
from modules.utils.logger import log_tool_call, log_agent_action, log_error


class CreatePPTTool(BaseTool):
    """Generate a PowerPoint presentation."""

    name = "create_ppt"
    description = "Generate a PowerPoint (.pptx) file with slides, backgrounds, and content."
    category = "ppt"
    parameters = {
        "type": "object",
        "properties": {
            "config": {
                "type": "object",
                "description": "PPT configuration including slides, color_scheme, etc.",
                "properties": {
                    "slides": {
                        "type": "array",
                        "description": "List of slide definitions",
                        "items": {
                            "type": "object",
                            "properties": {
                                "type": {"type": "string", "description": "Slide type: title_slide, section_header, content"},
                                "title": {"type": "string"},
                                "subtitle": {"type": "string"},
                                "content": {"type": "array", "items": {"type": "string"}},
                                "notes": {"type": "string"},
                                "layout": {"type": "string"},
                                "background": {"type": "object"}
                            }
                        }
                    },
                    "color_scheme": {
                        "type": "string",
                        "description": "Color scheme: modern_blue, corporate_gray, elegant_green, dark_elegant, tech",
                        "default": "modern_blue"
                    }
                }
            }
        },
        "required": ["config"]
    }

    def execute(self, config: Dict[str, Any] = None, **kwargs) -> str:
        config = config or {}
        log_tool_call("create_ppt()")
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE

            output_dir = os.path.join(
                os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output"
            )
            os.makedirs(output_dir, exist_ok=True)

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            slides_data = config.get("slides", [])
            color_scheme = config.get("color_scheme", "modern_blue")

            schemes = {
                "modern_blue": {
                    "primary": RGBColor(0x1A, 0x3C, 0x6E),
                    "secondary": RGBColor(0x2D, 0x7D, 0xD2),
                    "accent": RGBColor(0xFF, 0x8C, 0x00),
                    "light": RGBColor(0xF0, 0xF4, 0xFA),
                    "dark": RGBColor(0x0D, 0x1B, 0x2A),
                    "text": RGBColor(0x33, 0x33, 0x33),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "corporate_gray": {
                    "primary": RGBColor(0x2C, 0x3E, 0x50),
                    "secondary": RGBColor(0x7F, 0x8C, 0x8D),
                    "accent": RGBColor(0x34, 0x98, 0xDB),
                    "light": RGBColor(0xEC, 0xF0, 0xF1),
                    "dark": RGBColor(0x1A, 0x25, 0x2F),
                    "text": RGBColor(0x2C, 0x3E, 0x50),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "elegant_green": {
                    "primary": RGBColor(0x1E, 0x56, 0x36),
                    "secondary": RGBColor(0x4C, 0xAF, 0x50),
                    "accent": RGBColor(0xFF, 0xD7, 0x00),
                    "light": RGBColor(0xE8, 0xF5, 0xE9),
                    "dark": RGBColor(0x0D, 0x2B, 0x1A),
                    "text": RGBColor(0x33, 0x33, 0x33),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "dark_elegant": {
                    "primary": RGBColor(0x0D, 0x0D, 0x0D),
                    "secondary": RGBColor(0x1A, 0x1A, 0x2E),
                    "accent": RGBColor(0xC9, 0xA9, 0x4C),
                    "light": RGBColor(0x2A, 0x2A, 0x3E),
                    "dark": RGBColor(0x00, 0x00, 0x00),
                    "text": RGBColor(0xE0, 0xE0, 0xE0),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
                "tech": {
                    "primary": RGBColor(0x0A, 0x0A, 0x1A),
                    "secondary": RGBColor(0x00, 0xD4, 0xFF),
                    "accent": RGBColor(0xBB, 0x86, 0xFC),
                    "light": RGBColor(0x1A, 0x1A, 0x3A),
                    "dark": RGBColor(0x00, 0x00, 0x00),
                    "text": RGBColor(0xE0, 0xE0, 0xE0),
                    "white": RGBColor(0xFF, 0xFF, 0xFF),
                },
            }

            scheme = schemes.get(color_scheme, schemes["modern_blue"])

            for slide_data in slides_data:
                slide_type = slide_data.get("type", "content")
                title = slide_data.get("title", "")
                content = slide_data.get("content", [])
                bg_config = slide_data.get("background", {})

                if slide_type == "title_slide":
                    layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(layout)
                    self._apply_background(slide, bg_config, scheme)

                    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(44)
                    p.font.color.rgb = scheme["white"] if "dark" in color_scheme or "tech" in color_scheme else scheme["primary"]
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                    subtitle = slide_data.get("subtitle", "")
                    if subtitle:
                        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.333), Inches(1))
                        tf2 = txBox2.text_frame
                        p2 = tf2.paragraphs[0]
                        p2.text = subtitle
                        p2.font.size = Pt(24)
                        p2.font.color.rgb = scheme["accent"]
                        p2.alignment = PP_ALIGN.CENTER

                elif slide_type == "section_header":
                    layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(layout)
                    self._apply_background(slide, bg_config, scheme)

                    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.5))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(40)
                    p.font.color.rgb = scheme["accent"]
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

                else:
                    layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(layout)
                    self._apply_background(slide, bg_config, scheme)

                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = scheme["primary"]
                    shape.line.fill.background()

                    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.9))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(32)
                    p.font.color.rgb = scheme["white"]
                    p.font.bold = True

                    if content:
                        left_margin = 0.8 if slide_data.get("layout") != "two_column" else 0.5
                        txBox2 = slide.shapes.add_textbox(
                            Inches(left_margin), Inches(1.5),
                            Inches(11.733 if slide_data.get("layout") != "two_column" else 5.5),
                            Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.word_wrap = True

                        for i, item in enumerate(content):
                            if i == 0:
                                p = tf2.paragraphs[0]
                            else:
                                p = tf2.add_paragraph()
                            p.text = item
                            p.font.size = Pt(18)
                            p.font.color.rgb = scheme["text"]
                            p.space_after = Pt(8)
                            p.level = 0
                            if item.startswith("  "):
                                p.level = 1
                                p.font.size = Pt(16)

                notes = slide_data.get("notes", "")
                if notes and slide.has_notes_slide:
                    slide.notes_slide.notes_text_frame.text = notes

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
            filepath = os.path.join(output_dir, filename)
            prs.save(filepath)
            log_agent_action(f"PPT saved: {filepath}")
            return filepath

        except Exception as e:
            log_error(f"PPT generation failed: {e}")
            raise

    @staticmethod
    def _apply_background(slide, bg_config: Dict[str, Any], scheme: Dict):
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
        from lxml import etree

        bg_type = bg_config.get("type", "solid")
        color_1 = bg_config.get("color_1")
        color_2 = bg_config.get("color_2")

        def _to_rgb(c):
            if isinstance(c, list) and len(c) >= 3:
                return RGBColor(*c[:3])
            return scheme.get("light", RGBColor(240, 244, 250))

        c1 = _to_rgb(color_1) if color_1 else scheme.get("primary", RGBColor(26, 60, 110))
        c2 = _to_rgb(color_2) if color_2 else scheme.get("secondary", RGBColor(45, 125, 210))

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = c1

        if bg_type == "gradient":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = c2
            shape.line.fill.background()

            spPr = shape._element.find(qn('p:spPr'))
            if spPr is not None:
                solidFill_elem = spPr.find(qn('a:solidFill'))
                if solidFill_elem is not None:
                    srgbClr = solidFill_elem.find(qn('a:srgbClr'))
                    if srgbClr is not None:
                        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                        alpha.set('val', '50000')
